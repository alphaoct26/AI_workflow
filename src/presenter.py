import sqlite3
import pandas as pd
import matplotlib.pyplot as plt
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from src.config import (
    THEME_COLORS, THEME_HEX, FONT_TITLE, FONT_BODY, DB_PATH, CHART_PATH, PPTX_PATH
)
from src.analyst import BusinessInsights

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger("Presenter")

def generate_matplotlib_chart(workspace_id: str = "default"):
    """Queries the Gold summary and creates a clean, premium bar chart of Monthly Revenue."""
    logger.info("Presenter: Creating Matplotlib visual chart...")
    
    from src.db_adapter import DatabaseAdapter
    db = DatabaseAdapter(workspace_id)
    conn = db.get_connection()
    df = pd.read_sql_query("SELECT * FROM gold_monthly_metrics", conn)
    conn.close()
    
    if df.empty:
        logger.warning("Presenter: Gold table is empty. Skipping chart generation.")
        return
        
    # Dynamically resolve X (grouping) and Y (numeric) columns
    x_col = df.columns[0]
    y_col = None
    for col in df.columns[1:]:
        if pd.api.types.is_numeric_dtype(df[col]):
            y_col = col
            break
    if not y_col:
        y_col = df.columns[1] if len(df.columns) > 1 else df.columns[0]
        
    # Sort by X column
    df = df.sort_values(by=x_col, ascending=True)
    
    # Enable modern style
    plt.rcParams['font.family'] = FONT_TITLE
    fig, ax = plt.subplots(figsize=(7, 5), dpi=300)
    fig.patch.set_facecolor(THEME_HEX['light_bg'])
    ax.set_facecolor(THEME_HEX['light_bg'])
    
    max_val = df[y_col].max()
    should_scale = max_val > 1000
    y_vals = df[y_col] / 1000.0 if should_scale else df[y_col]
    
    is_currency = any(kw in y_col.lower() for kw in ("revenue", "cost", "salary", "amount", "price", "copay", "fee"))
    unit_suffix = "k" if should_scale else ""
    prefix = "$" if is_currency else ""
    
    # Draw bars with Indigo primary color and clean border rounded aesthetic
    bars = ax.bar(
        df[x_col].astype(str), 
        y_vals, 
        color=THEME_HEX['primary'], 
        width=0.55, 
        edgecolor=THEME_HEX['primary'],
        alpha=0.9
    )
    
    # Add values on top of the bars
    for bar in bars:
        yval = bar.get_height()
        label_text = f"{prefix}{yval:,.1f}{unit_suffix}" if should_scale else f"{prefix}{yval:,.0f}"
        ax.text(
            bar.get_x() + bar.get_width()/2.0, 
            yval, 
            label_text, 
            ha='center', 
            va='bottom', 
            fontsize=9, 
            color=THEME_HEX['text_muted'],
            fontweight='bold'
        )
        
    # Styling labels and titles
    title_text = f"MONTHLY {y_col.upper()} OVERVIEW"
    y_label_text = f"{y_col} (in thousands)" if should_scale else y_col
    if is_currency:
        y_label_text += " USD" if should_scale else " ($)"
        
    ax.set_title(title_text, fontsize=12, fontweight='bold', pad=20, color=THEME_HEX['text_dark'])
    ax.set_ylabel(y_label_text, fontsize=10, fontweight='bold', color=THEME_HEX['text_muted'])
    ax.set_xlabel(x_col, fontsize=10, fontweight='bold', color=THEME_HEX['text_muted'])
    
    # Adjust ticks
    ax.tick_params(colors=THEME_HEX['text_muted'], labelsize=9)
    
    # Remove outer spines for borderless premium feel
    for spine in ['top', 'right', 'left', 'bottom']:
        ax.spines[spine].set_visible(False)
        
    # Add thin light gridlines for y-axis
    ax.grid(axis='y', linestyle='--', alpha=0.6, color=THEME_HEX['grid_color'])
    ax.set_axisbelow(True) # Keep grid behind bars
    
    plt.tight_layout()
    from src.config import get_workspace_file_paths
    paths = get_workspace_file_paths(workspace_id)
    plt.savefig(paths["chart_png"], facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close()
    
    logger.info(f"Presenter: Chart generated and saved to {paths['chart_png']}.")

def apply_text_formatting(paragraph, text, size_pt, color_rgb, font_name, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    """Utility to safely format paragraphs in python-pptx."""
    paragraph.text = text
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size_pt)
    paragraph.font.color.rgb = RGBColor(*color_rgb)
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.alignment = alignment

def add_slide_card(slide, left, top, width, height):
    """Draws a white rounded panel box (dashboard card style) to back the text."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*THEME_COLORS['card_bg'])
    card.line.color.rgb = RGBColor(226, 232, 240) # Slate 200 light border
    card.line.width = Pt(1)
    return card

def create_presentation_deck(insights: BusinessInsights, workspace_id: str = "default"):
    """Builds the final 3-slide e-commerce presentation deck using parsed JSON insights."""
    logger.info("Presenter: Compiling PowerPoint deck...")
    
    # Initialize standard presentation
    prs = Presentation()
    
    # Set to 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layouts (usually index 6 in default templates is blank)
    blank_layout = prs.slide_layouts[6]
    
    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    
    # Set dark background fill
    bg_1 = slide_1.background
    fill_1 = bg_1.fill
    fill_1.solid()
    fill_1.fore_color.rgb = RGBColor(*THEME_COLORS['dark_bg'])
    
    # Decorative accent line
    accent_bar = slide_1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.15), Inches(3.2)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(*THEME_COLORS['secondary'])
    accent_bar.line.fill.background()
    
    # Title Text Frame (Single text frame to prevent overlapping title/subtitle)
    title_box = slide_1.shapes.add_textbox(Inches(1.8), Inches(2.0), Inches(10), Inches(3.5))
    tf_1 = title_box.text_frame
    tf_1.word_wrap = True
    
    # Title Paragraph
    p_title = tf_1.paragraphs[0]
    apply_text_formatting(p_title, "AUTO-ANALYST REPORT", 40, THEME_COLORS['light_bg'], FONT_TITLE, bold=True)
    p_title.space_after = Pt(8)
    
    # Subtitle Paragraph
    p_sub = tf_1.add_paragraph()
    apply_text_formatting(p_sub, f"Dynamic Business Performance & Insights  |  Workspace: {workspace_id}", 16, THEME_COLORS['text_muted'], FONT_BODY, italic=True)
    p_sub.space_after = Pt(28)
    
    # Metadata Paragraph
    p_meta = tf_1.add_paragraph()
    run_date = datetime.now().strftime("%B %d, %Y")
    apply_text_formatting(p_meta, f"Generated automatically on {run_date}\nPowered by Gemini & DuckDB OLAP Engine", 11, THEME_COLORS['text_muted'], FONT_BODY)
    
    # ==========================================
    # SLIDE 2: Executive Summary & Key Findings (Light Theme)
    # ==========================================
    slide_2 = prs.slides.add_slide(blank_layout)
    bg_2 = slide_2.background
    fill_2 = bg_2.fill
    fill_2.solid()
    fill_2.fore_color.rgb = RGBColor(*THEME_COLORS['light_bg'])
    
    # Header
    header_2 = slide_2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8))
    apply_text_formatting(
        header_2.text_frame.paragraphs[0], "Executive Summary & Performance Audits", 26, THEME_COLORS['text_dark'], FONT_TITLE, bold=True
    )
    
    # Column 1: Executive Summary Card (Left)
    col_width = Inches(5.6)
    card_height = Inches(5.1)
    left_col_1 = Inches(0.8)
    add_slide_card(slide_2, left_col_1, Inches(1.5), col_width, card_height)
    
    summary_box = slide_2.shapes.add_textbox(left_col_1 + Inches(0.2), Inches(1.7), col_width - Inches(0.4), card_height - Inches(0.4))
    tf_summary = summary_box.text_frame
    tf_summary.word_wrap = True
    
    p_sum_head = tf_summary.paragraphs[0]
    apply_text_formatting(p_sum_head, "BUSINESS EXECUTIVE OVERVIEW", 14, THEME_COLORS['primary'], FONT_TITLE, bold=True)
    p_sum_head.space_after = Pt(14)
    
    p_sum_body = tf_summary.add_paragraph()
    apply_text_formatting(p_sum_body, insights.executive_summary, 13, THEME_COLORS['text_dark'], FONT_BODY)
    p_sum_body.line_spacing = 1.2
    
    # Column 2: Key Findings Card (Right)
    left_col_2 = Inches(6.9)
    add_slide_card(slide_2, left_col_2, Inches(1.5), col_width, card_height)
    
    findings_box = slide_2.shapes.add_textbox(left_col_2 + Inches(0.2), Inches(1.7), col_width - Inches(0.4), card_height - Inches(0.4))
    tf_findings = findings_box.text_frame
    tf_findings.word_wrap = True
    
    p_find_head = tf_findings.paragraphs[0]
    apply_text_formatting(p_find_head, "CORE BUSINESS FINDINGS & ANOMALIES", 14, THEME_COLORS['primary'], FONT_TITLE, bold=True)
    p_find_head.space_after = Pt(14)
    
    for finding in insights.key_findings:
        p_item = tf_findings.add_paragraph()
        p_item.space_after = Pt(12)
        p_item.level = 0
        p_item.line_spacing = 1.15
        
        # Format bullet point
        apply_text_formatting(p_item, f"▪  {finding}", 12, THEME_COLORS['text_dark'], FONT_BODY)
        
    # ==========================================
    # SLIDE 3: Charts & Recommendations (Light Theme)
    # ==========================================
    slide_3 = prs.slides.add_slide(blank_layout)
    bg_3 = slide_3.background
    fill_3 = bg_3.fill
    fill_3.solid()
    fill_3.fore_color.rgb = RGBColor(*THEME_COLORS['light_bg'])
    
    # Title
    header_3 = slide_3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8))
    apply_text_formatting(
        header_3.text_frame.paragraphs[0], "Financial Trends & AI Recommendations", 26, THEME_COLORS['text_dark'], FONT_TITLE, bold=True
    )
    
    # Column 1: Matplotlib Chart Card (Left)
    add_slide_card(slide_3, Inches(0.8), Inches(1.5), col_width, card_height)
    
    # Insert visual image inside the card
    from src.config import get_workspace_file_paths
    paths = get_workspace_file_paths(workspace_id)
    slide_3.shapes.add_picture(
        str(paths["chart_png"]), 
        Inches(1.0), 
        Inches(1.75), 
        width=col_width - Inches(0.4), 
        height=card_height - Inches(0.5)
    )
    
    # Column 2: Actionable Recommendations Card (Right)
    add_slide_card(slide_3, left_col_2, Inches(1.5), col_width, card_height)
    rec_box = slide_3.shapes.add_textbox(left_col_2 + Inches(0.2), Inches(1.7), col_width - Inches(0.4), card_height - Inches(0.4))
    tf_rec = rec_box.text_frame
    tf_rec.word_wrap = True
    
    p_rec_head = tf_rec.paragraphs[0]
    apply_text_formatting(p_rec_head, "RECOMMENDED STRATEGIC ACTIONS", 14, THEME_COLORS['secondary'], FONT_TITLE, bold=True)
    p_rec_head.space_after = Pt(14)
    
    for action in insights.recommended_actions:
        p_action = tf_rec.add_paragraph()
        p_action.space_after = Pt(12)
        p_action.level = 0
        p_action.line_spacing = 1.15
        
        # Format bullet point
        apply_text_formatting(p_action, f"✔  {action}", 12, THEME_COLORS['text_dark'], FONT_BODY)
        
    # Save PPTX
    prs.save(str(paths["pptx_report"]))
    logger.info(f"Presenter: Presentation saved successfully to {paths['pptx_report']}.")
